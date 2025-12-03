from pptx import Presentation

prs = Presentation()
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Hello, World!"
slide.placeholders[1].text = "Python-pptx test"

prs.save('test.pptx')
print("Saved test.pptx")
